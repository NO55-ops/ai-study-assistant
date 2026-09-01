from fastapi import FastAPI, APIRouter, File, UploadFile, Request, HTTPException, Response, Header, Query
from fastapi.responses import StreamingResponse
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware

from typing import List, Optional, Dict, Any

try:
    from supabase import create_client, Client
except ModuleNotFoundError:
    create_client = None
    Client = Any

import os
import logging
import certifi
from urllib.parse import urlparse, quote_plus
from pathlib import Path
from pydantic import BaseModel, Field, EmailStr
import uuid
from datetime import datetime, timezone, timedelta
import bcrypt
import jwt
import requests
import io
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from openai import AsyncOpenAI
import json

ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

SUPABASE_URL = os.environ.get("SUPABASE_URL", "")
SUPABASE_KEY = os.environ.get("SUPABASE_KEY", "")

MEMORY_STORE: Dict[str, List[Dict[str, Any]]] = {}


def normalize_record(record: Optional[Dict[str, Any]]) -> Optional[Dict[str, Any]]:
    if record is None:
        return None
    normalized = dict(record)
    if "id" not in normalized and "_id" in normalized:
        normalized["id"] = normalized["_id"]
    if "_id" not in normalized and "id" in normalized:
        normalized["_id"] = normalized["id"]
    return normalized


def normalize_supabase_id(record: Optional[Dict[str, Any]]) -> Optional[Dict[str, Any]]:
    if record is None:
        return None
    normalized = normalize_record(record)
    if normalized is None:
        return None
    if "id" not in normalized and "_id" in normalized:
        normalized["id"] = normalized["_id"]
    normalized.pop("_id", None)
    return normalized


def ensure_response_id(record: Optional[Dict[str, Any]]) -> Optional[Dict[str, Any]]:
    normalized = normalize_supabase_id(record)
    if normalized is None:
        return None
    if "id" not in normalized and "_id" in normalized:
        normalized["id"] = normalized["_id"]
    normalized.pop("_id", None)
    if "id" not in normalized:
        normalized["id"] = str(uuid.uuid4())
    return normalized


class MemoryResult:
    def __init__(self, data: Optional[List[Dict[str, Any]]] = None, count: Optional[int] = None):
        self.data = data or []
        self.count = count


class OperationResult:
    def __init__(self, matched_count: int = 0, modified_count: int = 0, deleted_count: int = 0):
        self.matched_count = matched_count
        self.modified_count = modified_count
        self.deleted_count = deleted_count


class MemoryTable:
    def __init__(self, name: str):
        self.name = name

    def select(self, *fields, **kwargs):
        return MemoryQueryBuilder(self.name, fields=fields, count=kwargs.get("count"))

    def insert(self, document: Dict[str, Any]):
        return MemoryQueryBuilder(self.name).insert(document)

    def update(self, payload: Dict[str, Any]):
        return MemoryQueryBuilder(self.name).update(payload)

    def delete(self):
        return MemoryQueryBuilder(self.name).delete()


class MemoryQueryBuilder:
    def __init__(self, table_name: str, fields: Optional[tuple] = None, count: Optional[str] = None):
        self.table_name = table_name
        self.fields = fields or ("*",)
        self.count = count
        self.filters = []
        self.sort_spec = None
        self.limit_value = None
        self.payload = None
        self.action = "read"

    def eq(self, field: str, value: Any):
        self.filters.append(
            lambda row, f=field, v=value: (row.get(f) if row.get(f) is not None else row.get("id" if f == "_id" else "_id")) == v
        )
        return self

    def in_(self, field: str, values: List[Any]):
        self.filters.append(
            lambda row, f=field, v=set(values): (row.get(f) if row.get(f) is not None else row.get("id" if f == "_id" else "_id")) in v
        )
        return self

    def order(self, field: str, desc: bool = False):
        self.sort_spec = (field, desc)
        return self

    def limit(self, value: int):
        self.limit_value = value
        return self

    def insert(self, document: Dict[str, Any]):
        self.action = "insert"
        self.payload = dict(document)
        return self

    def update(self, payload: Dict[str, Any]):
        self.action = "update"
        self.payload = payload
        return self

    def delete(self):
        self.action = "delete"
        return self

    def execute(self):
        rows = MEMORY_STORE.setdefault(self.table_name, [])
        if self.action == "insert":
            payload = dict(self.payload or {})
            if "id" not in payload and "_id" not in payload:
                payload["id"] = str(uuid.uuid4())
            rows.append(dict(payload))
            return MemoryResult(data=[dict(payload)])

        filtered: List[Dict[str, Any]] = []
        for row in rows:
            record = normalize_record(dict(row))
            if record is not None and all(f(record) for f in self.filters):
                filtered.append(record)

        if self.action == "delete":
            for row in list(rows):
                record = normalize_record(dict(row))
                if record is not None and all(f(record) for f in self.filters):
                    rows.remove(row)
            return MemoryResult(data=[])

        if self.action == "update":
            updated_rows: List[Dict[str, Any]] = []
            payload = dict(self.payload or {})
            for row in rows:
                record = normalize_record(dict(row))
                if record is not None and all(f(record) for f in self.filters):
                    record.update(payload)
                    updated_rows.append(record)
                    rows[rows.index(row)] = record
            return MemoryResult(data=updated_rows)

        if self.sort_spec:
            field, desc = self.sort_spec
            filtered.sort(key=lambda row: (row.get(field) is None, str(row.get(field, ""))), reverse=desc)

        if self.limit_value is not None:
            filtered = filtered[: self.limit_value]

        if self.count == "exact":
            return MemoryResult(data=filtered, count=len(filtered))

        if self.fields and self.fields != ("*",):
            projected = []
            for row in filtered:
                item = {}
                for field in self.fields:
                    if field == "id":
                        item["id"] = row.get("id", row.get("_id"))
                    elif field == "*":
                        item = dict(row)
                    else:
                        item[field] = row.get(field)
                projected.append(item)
            return MemoryResult(data=projected)

        return MemoryResult(data=filtered)


class MemorySupabaseClient:
    def table(self, name: str):
        return MemoryTable(name)


logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

supabase: Any = None
if create_client is not None and SUPABASE_URL and SUPABASE_KEY:
    supabase = create_client(SUPABASE_URL, SUPABASE_KEY)
else:
    logger.warning("Supabase client unavailable or credentials not configured; falling back to local in-memory storage for development.")
    supabase = MemorySupabaseClient()


class SupabaseQuery:
    def __init__(self, collection_name: str, query: Optional[Dict[str, Any]] = None, projection: Optional[Dict[str, Any]] = None):
        self.collection_name = collection_name
        self.query = query or {}
        self.projection = projection or {}
        self._sort_spec = None
        self._limit_value = None

    def sort(self, *args):
        if len(args) == 1 and isinstance(args[0], list):
            self._sort_spec = args[0]
        elif len(args) == 2:
            self._sort_spec = [(args[0], -1 if args[1] < 0 else 1)]
        else:
            self._sort_spec = None
        return self

    def limit(self, count: int):
        self._limit_value = count
        return self

    @staticmethod
    def _normalize_key(key: str) -> str:
        return "id" if key in ("_id", "id") else key

    @staticmethod
    def _normalize_filter(query: Dict[str, Any]) -> Dict[str, Any]:
        normalized = {}
        for key, value in (query or {}).items():
            normalized[SupabaseQuery._normalize_key(key)] = value
        return normalized

    @staticmethod
    def _apply_projection(document: Dict[str, Any], projection: Dict[str, Any]) -> Dict[str, Any]:
        if not document or not isinstance(projection, dict):
            return document
        projected = dict(document)
        for key, value in projection.items():
            if value == 0:
                projected.pop(key, None)
                projected.pop(key.replace("_id", "id"), None)
        return projected

    @staticmethod
    def _build_table(name: str):
        return supabase.table(name)

    async def _execute(self):
        table = self._build_table(self.collection_name).select("*")
        for field, value in self._normalize_filter(self.query).items():
            if isinstance(value, dict):
                continue
            if isinstance(value, (list, tuple, set)):
                table = table.in_(field, list(value))
            else:
                table = table.eq(field, value)

        if self._sort_spec:
            for field, direction in self._sort_spec:
                table = table.order(field, desc=(direction < 0))

        if self._limit_value is not None:
            table = table.limit(self._limit_value)

        result = table.execute()
        rows = result.data or []
        return [self._apply_projection(row, self.projection) for row in rows]

    async def to_list(self, fallback_limit: Optional[int] = None):
        if fallback_limit is not None and self._limit_value is None:
            self._limit_value = fallback_limit
        return await self._execute()


class SupabaseCollection:
    def __init__(self, name: str):
        self.name = name

    def find(self, query: Optional[Dict[str, Any]] = None, projection: Optional[Dict[str, Any]] = None):
        return SupabaseQuery(self.name, query=query, projection=projection)

    async def find_one(self, query: Optional[Dict[str, Any]] = None, projection: Optional[Dict[str, Any]] = None, sort=None):
        q = SupabaseQuery(self.name, query=query, projection=projection)
        if sort is not None:
            q.sort(sort)
        rows = await q.to_list(1)
        return rows[0] if rows else None

    async def insert_one(self, document: Dict[str, Any]):
        table = SupabaseQuery._build_table(self.name)
        payload = normalize_supabase_id(document) or dict(document)
        payload.pop("_id", None)
        result = table.insert(payload).execute()
        data = result.data or []
        return data[0] if data else None

    async def update_one(self, query: Dict[str, Any], update: Dict[str, Any]):
        update_payload = update.get("$set", update) if isinstance(update, dict) and "$set" in update else update
        table = SupabaseQuery._build_table(self.name).update(update_payload)
        for field, value in SupabaseQuery._normalize_filter(query).items():
            table = table.eq(field, value)
        result = table.execute()
        rows = result.data or []
        return OperationResult(matched_count=1 if rows else 0, modified_count=1 if rows else 0)

    async def delete_one(self, query: Dict[str, Any]):
        table = SupabaseQuery._build_table(self.name).delete()
        for field, value in SupabaseQuery._normalize_filter(query).items():
            table = table.eq(field, value)
        result = table.execute()
        rows = result.data or []
        return OperationResult(deleted_count=1 if rows else 0)

    async def count_documents(self, query: Optional[Dict[str, Any]] = None):
        result = SupabaseQuery._build_table(self.name).select("id", count="exact")
        for field, value in SupabaseQuery._normalize_filter(query or {}).items():
            if isinstance(value, (list, tuple, set)):
                result = result.in_(field, list(value))
            else:
                result = result.eq(field, value)
        data = result.execute()
        if hasattr(data, "count") and data.count is not None:
            return data.count
        return len(data.data or [])

    async def create_index(self, field: str, unique: bool = False):
        return None


class SupabaseDatabase:
    users = SupabaseCollection("users")
    documents = SupabaseCollection("documents")
    summaries = SupabaseCollection("summaries")
    quizzes = SupabaseCollection("quizzes")
    flashcard_sets = SupabaseCollection("flashcard_sets")
    study_plans = SupabaseCollection("study_plans")
    streaks = SupabaseCollection("streaks")
    notes = SupabaseCollection("notes")


db = SupabaseDatabase()

app = FastAPI()
api_router = APIRouter(prefix="/api")

# Storage setup (local filesystem)
STORAGE_DIR = Path(os.environ.get("STORAGE_DIR", str(ROOT_DIR / "storage")))
STORAGE_DIR.mkdir(parents=True, exist_ok=True)
APP_NAME = "ai-study-assistant"

# AI provider configuration
OPENROUTER_API_KEY = os.environ.get("OPENROUTER_API_KEY")
OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY")
OPENAI_BASE_URL = os.environ.get("OPENAI_BASE_URL")

if OPENROUTER_API_KEY:
    OPENAI_API_KEY = OPENROUTER_API_KEY
    OPENAI_BASE_URL = "https://openrouter.ai/api/v1"
    OPENAI_MODEL = os.environ.get("OPENAI_MODEL", "openrouter/free")
elif OPENAI_BASE_URL:
    OPENAI_MODEL = os.environ.get("OPENAI_MODEL", os.environ.get("OLLAMA_MODEL", "gemma4:12b"))
else:
    OPENAI_API_KEY = OPENAI_API_KEY or "ollama"
    OPENAI_BASE_URL = "http://localhost:11434/v1"
    OPENAI_MODEL = os.environ.get("OPENAI_MODEL", os.environ.get("OLLAMA_MODEL", "gemma4:12b"))

openai_client = AsyncOpenAI(
    base_url=OPENAI_BASE_URL,
    api_key=OPENAI_API_KEY,
)


def put_object(path: str, data: bytes, content_type: str) -> dict:
    """Save file to local filesystem storage."""
    full_path = STORAGE_DIR / path
    full_path.parent.mkdir(parents=True, exist_ok=True)
    full_path.write_bytes(data)
    return {"path": path, "size": len(data), "content_type": content_type}


def get_object(path: str) -> tuple:
    """Read file from local filesystem storage."""
    full_path = STORAGE_DIR / path
    if not full_path.exists():
        raise FileNotFoundError(f"Object not found: {path}")
    return full_path.read_bytes(), "application/octet-stream"


async def llm_generate(system: str, prompt: str, model: Optional[str] = None) -> str:
    """Generate a complete LLM response (non-streaming)."""
    if not openai_client:
        raise HTTPException(
    status_code=503,
    detail="Ollama is not running or not configured."
)
    english_system = (
        "You are a clear, helpful study assistant. Always answer in standard English only. "
        "Keep your response concise, well-structured, and easy for students to understand."
    )
    resp = await openai_client.chat.completions.create(
       model=model or OPENAI_MODEL,
        messages=[
            {"role": "system", "content": english_system + " " + system},
            {"role": "user", "content": prompt},
        ],
    )
    return resp.choices[0].message.content or ""


async def llm_stream(system: str, prompt: str, model: Optional[str] = None):
    """Stream LLM response as async iterator of text chunks."""
    if not openai_client:
        raise HTTPException(
            status_code=503,
            detail="Ollama is not running or not configured.",
        )

    english_system = (
        "You are a clear, helpful study assistant. Always answer in standard English only. "
        "Keep your response concise, well-structured, and easy for students to understand."
    )

    stream = await openai_client.chat.completions.create(
        model=model or OPENAI_MODEL,
        messages=[
            {"role": "system", "content": english_system + " " + system},
            {"role": "user", "content": prompt},
        ],
        stream=True,
    )

    async for chunk in stream:
        delta = chunk.choices[0].delta.content if chunk.choices else None
        if delta:
            yield delta
# Auth utilities
JWT_ALGORITHM = "HS256"

def get_jwt_secret() -> str:
    return os.environ.get("JWT_SECRET", "dev-secret-key")

def hash_password(password: str) -> str:
    salt = bcrypt.gensalt()
    hashed = bcrypt.hashpw(password.encode("utf-8"), salt)
    return hashed.decode("utf-8")

def verify_password(plain_password: str, hashed_password: str) -> bool:
    return bcrypt.checkpw(plain_password.encode("utf-8"), hashed_password.encode("utf-8"))

def create_access_token(user_id: str, email: str) -> str:
    payload = {"sub": user_id, "email": email, "exp": datetime.now(timezone.utc) + timedelta(minutes=15), "type": "access"}
    return jwt.encode(payload, get_jwt_secret(), algorithm=JWT_ALGORITHM)

def create_refresh_token(user_id: str) -> str:
    payload = {"sub": user_id, "exp": datetime.now(timezone.utc) + timedelta(days=7), "type": "refresh"}
    return jwt.encode(payload, get_jwt_secret(), algorithm=JWT_ALGORITHM)

async def get_current_user(request: Request) -> dict:
    token = request.cookies.get("access_token")

    if not token:
        auth_header = request.headers.get("Authorization", "")
        if auth_header.startswith("Bearer "):
            token = auth_header[7:]

    if not token:
        raise HTTPException(status_code=401, detail="Not authenticated")

    try:
        payload = jwt.decode(
            token,
            get_jwt_secret(),
            algorithms=[JWT_ALGORITHM]
        )

        if payload.get("type") != "access":
            raise HTTPException(
                status_code=401,
                detail="Invalid token type"
            )

        user = await db.users.find_one({"id": payload["sub"]})
        if not user:
            user = await db.users.find_one({"_id": payload["sub"]})

        user = normalize_record(user)

        if not user:
            raise HTTPException(
                status_code=401,
                detail="User not found"
            )

        user.pop("password_hash", None)
        return user

    except jwt.ExpiredSignatureError:
        raise HTTPException(
            status_code=401,
            detail="Token expired"
        )

    except jwt.InvalidTokenError:
        raise HTTPException(
            status_code=401,
            detail="Invalid token"
        )


# Models
class RegisterRequest(BaseModel):
    email: EmailStr
    password: str
    name: str


class LoginRequest(BaseModel):
    email: EmailStr
    password: str


class OnboardingRequest(BaseModel):
    grade_level: str
    subjects: List[str]
    curriculum: str
    exam_dates: Optional[List[Dict[str, str]]] = []
    daily_study_goal: int = 60


class SummaryRequest(BaseModel):
    document_id: str
    summary_type: str = "detailed"


class TutorRequest(BaseModel):
    document_id: str
    question: str
    difficulty_mode: str = "intermediate"
    session_id: Optional[str] = None


class QuizRequest(BaseModel):
    document_ids: List[str]
    question_types: List[str]
    difficulty: str = "medium"
    num_questions: int = 10


class FlashcardRequest(BaseModel):
    document_id: str


class StudyPlanRequest(BaseModel):
    exam_dates: List[Dict[str, str]]
    daily_study_time: int


class NoteRequest(BaseModel):
    title: str
    content: str = ""
    tags: List[str] = []
    color: str = "butter"
    is_pinned: bool = False
    document_id: Optional[str] = None


# Helper: Extract text from files
def extract_text_from_file(file_content: bytes, file_type: str) -> str:
    try:
        if file_type == "application/pdf":
            pdf_reader = PdfReader(io.BytesIO(file_content))
            text = ""

            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"

            return text

        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = DocxDocument(io.BytesIO(file_content))

            return "\n".join(
                [para.text for para in doc.paragraphs]
            )

        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            prs = Presentation(io.BytesIO(file_content))
            text = ""

            for slide in prs.slides:
                for shape in slide.shapes:
                    txt = getattr(shape, "text", None)

                    if txt:
                        text += txt + "\n"

            return text

        elif file_type == "text/plain":
            return file_content.decode("utf-8")

        else:
            return ""

    except Exception as e:
        logger.error(f"Error extracting text: {e}")
        return ""


# Auth endpoints
@api_router.post("/auth/register")
async def register(req: RegisterRequest, response: Response):
    email = req.email.lower()

    result = (
        supabase
        .table("users")
        .select("*")
        .eq("email", email)
        .limit(1)
        .execute()
    )

    existing = result.data[0] if result.data else None

    if existing:
        raise HTTPException(
            status_code=400,
            detail="Email already registered"
        )

    user_id = str(uuid.uuid4())

    user_doc = {
        "id": user_id,
        "email": email,
        "name": req.name,
        "password_hash": hash_password(req.password),
        "onboarded": False,
        "created_at": datetime.now(timezone.utc).isoformat()
    }

    supabase.table("users").insert(user_doc).execute()

    access_token = create_access_token(user_id, email)
    refresh_token = create_refresh_token(user_id)

    response.set_cookie(
        key="access_token",
        value=access_token,
        httponly=True,
        secure=True,
        samesite="none",
        max_age=900,
        path="/"
    )

    response.set_cookie(
        key="refresh_token",
        value=refresh_token,
        httponly=True,
        secure=True,
        samesite="none",
        max_age=604800,
        path="/"
    )

    user_doc.pop("password_hash", None)

    return user_doc


@api_router.post("/auth/login")
async def login(req: LoginRequest, response: Response):
    email = req.email.lower()

    result = (
        supabase
        .table("users")
        .select("*")
        .eq("email", email)
        .limit(1)
        .execute()
    )

    user = result.data[0] if result.data else None
    user = normalize_record(user) if isinstance(user, dict) else user

    if not user:
        raise HTTPException(
            status_code=401,
            detail="Invalid email or password"
        )

    if not verify_password(
        req.password,
        str(user.get("password_hash", ""))
    ):
        raise HTTPException(
            status_code=401,
            detail="Invalid email or password"
        )

    user_id = str(user.get("id") or user.get("_id"))

    access_token = create_access_token(
        user_id,
        email
    )

    refresh_token = create_refresh_token(user_id)

    response.set_cookie(
        key="access_token",
        value=access_token,
        httponly=True,
        secure=True,
        samesite="none",
        max_age=900,
        path="/"
    )

    response.set_cookie(
        key="refresh_token",
        value=refresh_token,
        httponly=True,
        secure=True,
        samesite="none",
        max_age=604800,
        path="/"
    )

    user.pop("password_hash", None)
    user = normalize_record(user)
    return user


@api_router.post("/auth/logout")
async def logout(response: Response):
    response.delete_cookie(
        key="access_token",
        path="/"
    )

    response.delete_cookie(
        key="refresh_token",
        path="/"
    )

    return {
        "message": "Logged out successfully"
    }


@api_router.get("/auth/me")
async def get_me(request: Request):
    user = await get_current_user(request)

    return user


@api_router.post("/auth/onboarding")
async def complete_onboarding(
    req: OnboardingRequest,
    request: Request
):
    user = await get_current_user(request)

    update_data = {
        "school_level": req.grade_level,
        "study_preferences": {
            "subjects": req.subjects,
            "curriculum": req.curriculum,
            "exam_dates": req.exam_dates
        },
        "daily_study_goal": req.daily_study_goal,
        "onboarded": True
    }

    supabase.table("users").update(
        update_data
    ).eq(
        "id",
        user["id"]
    ).execute()

    existing_streak_result = (
        supabase
        .table("streaks")
        .select("*")
        .eq("user_id", user["id"])
        .limit(1)
        .execute()
    )

    existing_streak = (
        existing_streak_result.data[0]
        if existing_streak_result.data
        else None
    )

    if not existing_streak:

        streak_doc = {
            "id": str(uuid.uuid4()),
            "user_id": user["id"],
            "current_streak": 0,
            "longest_streak": 0,
            "last_activity_date": None,
            "total_study_minutes": 0
        }

        supabase.table(
            "streaks"
        ).insert(
            streak_doc
        ).execute()

    return {
        "message": "Onboarding completed"
    }

# File upload
@api_router.post("/documents/upload")
async def upload_document(request: Request, file: UploadFile = File(...)):
    user = await get_current_user(request)
    
    file_content = await file.read()
    filename = getattr(file, "filename", "") or ""
    ext = filename.split(".")[-1] if "." in filename else "bin"
    storage_path = f"{APP_NAME}/uploads/{user['id']}/{uuid.uuid4()}.{ext}"
    
    content_type = file.content_type or "application/octet-stream"
    result = put_object(storage_path, file_content, content_type)
    
    extracted_text = extract_text_from_file(file_content, content_type)
    
    doc_id = str(uuid.uuid4())
    doc = {
        "id": doc_id,
        "user_id": user["id"],
        "original_filename": file.filename,
        "storage_path": result["path"],
        "content_type": content_type,
        "size": result["size"],
        "extracted_text": extracted_text,
        "is_deleted": False,
        "created_at": datetime.now(timezone.utc).isoformat()
    }
    
    await db.documents.insert_one(doc)
    
    response_doc = normalize_supabase_id(doc) or dict(doc)
    response_doc.pop("extracted_text", None)
    return response_doc

@api_router.get("/documents")
async def list_documents(request: Request):
    user = await get_current_user(request)
    docs = await db.documents.find({"user_id": user["id"], "is_deleted": False}, {"extracted_text": 0}).to_list(1000)
    normalized_docs: List[Dict[str, Any]] = []
    for d in docs:
        normalized = normalize_supabase_id(d)
        if normalized is None:
            continue
        normalized.pop("_id", None)
        normalized_docs.append(normalized)
    return normalized_docs

@api_router.get("/documents/{doc_id}")
async def get_document(doc_id: str, request: Request):
    user = await get_current_user(request)
    doc = await db.documents.find_one({"_id": doc_id, "user_id": user["id"], "is_deleted": False})
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    doc = normalize_supabase_id(doc) or dict(doc)
    return doc

@api_router.delete("/documents/{doc_id}")
async def delete_document(doc_id: str, request: Request):
    user = await get_current_user(request)
    result = await db.documents.update_one({"_id": doc_id, "user_id": user["id"]}, {"$set": {"is_deleted": True}})
    if result.modified_count == 0:
        raise HTTPException(status_code=404, detail="Document not found")
    return {"message": "Document deleted"}

# AI Summary
@api_router.post("/ai/summary")
async def generate_summary(req: SummaryRequest, request: Request):
    user = await get_current_user(request)
    doc = await db.documents.find_one({"_id": req.document_id, "user_id": user["id"], "is_deleted": False})
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    
    existing = await db.summaries.find_one({"document_id": req.document_id, "summary_type": req.summary_type})
    if existing:
        existing = ensure_response_id(existing) or existing
        return existing
    
    text = doc.get("extracted_text", "")
    if not text:
        raise HTTPException(status_code=400, detail="No text content found in document")
    
    prompt_templates = {
        "short": f"Provide a concise summary (3-5 sentences) of the following study material:\n\n{text[:4000]}",
        "detailed": f"Create a comprehensive summary with:\n1. Main Topics\n2. Key Concepts\n3. Important Definitions\n4. Formulas (if any)\n5. Exam Tips\n\nMaterial:\n{text[:6000]}",
        "bullet_points": f"Summarize in clear bullet points:\n\n{text[:5000]}",
        "timeline": f"If applicable, create a timeline of events. Otherwise, organize by sequence:\n\n{text[:5000]}"
    }
    
    prompt = prompt_templates.get(req.summary_type, prompt_templates["detailed"])
    
    summary_content = await llm_generate(
        system="You are an expert study assistant. Create clear, well-organized summaries for students.",
        prompt=prompt,
    )
    
    summary_id = str(uuid.uuid4())
    summary_doc = {
        "id": summary_id,
        "document_id": req.document_id,
        "user_id": user["id"],
        "summary_type": req.summary_type,
        "content": summary_content,
        "created_at": datetime.now(timezone.utc).isoformat()
    }
    
    await db.summaries.insert_one(summary_doc)
    return normalize_supabase_id(summary_doc) or dict(summary_doc)

# AI Tutor (streaming)
@api_router.post("/ai/tutor")
async def ai_tutor(req: TutorRequest, request: Request):
    user = await get_current_user(request)
    doc = await db.documents.find_one({"_id": req.document_id, "user_id": user["id"], "is_deleted": False})
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    
    text = doc.get("extracted_text", "")[:6000]
    
    mode_prompts = {
        "beginner": "Explain concepts in very simple terms, use analogies, and avoid jargon.",
        "intermediate": "Provide clear explanations with examples. Balance simplicity with accuracy.",
        "advanced": "Provide in-depth explanations. You can use technical terminology.",
        "exam_prep": "Focus on exam-relevant information, common pitfalls, and test strategies."
    }
    
    system_msg = (
        "You are an English-speaking study tutor. Reply in standard English only. "
        "Use clear explanations, simple vocabulary, and concise structure. "
        f"{mode_prompts.get(req.difficulty_mode, mode_prompts['intermediate'])} "
        "Guide students with hints before giving direct answers. "
        f"Reference this material: {text}"
    )
    
    async def event_generator():
        async for delta in llm_stream(system=system_msg, prompt=req.question):
            yield f"data: {json.dumps({'content': delta})}\n\n"
        yield f"data: {json.dumps({'done': True})}\n\n"
    
    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"}
    )

# Quiz generation
@api_router.post("/ai/quiz")
async def generate_quiz(req: QuizRequest, request: Request):
    user = await get_current_user(request)

    safe_question_count = max(3, min(int(req.num_questions or 3), 8))
    max_context_chars = 3200
    combined_text = ""
    for doc_id in req.document_ids[:3]:
        doc = await db.documents.find_one({"_id": doc_id, "user_id": user["id"], "is_deleted": False})
        if doc:
            combined_text += (doc.get("extracted_text", "")[:1200] + "\n\n")

    combined_text = combined_text[:max_context_chars]

    if not combined_text:
        raise HTTPException(status_code=400, detail="No text found in documents")

    question_types_str = ", ".join(req.question_types[:3])

    prompt = f"""Generate exactly {safe_question_count} concise quiz questions from this material.
Question types: {question_types_str}
Difficulty: {req.difficulty}

Material:
{combined_text}

Return ONLY valid JSON in this exact shape:
{{
  "questions": [
    {{
      "type": "multiple_choice",
      "question": "...",
      "options": ["A", "B", "C", "D"],
      "correct_answer": "A",
      "explanation": "..."
    }}
  ]
}}"""

    quiz_content = await llm_generate(
        system="You are a fast quiz generator. Return valid JSON only, with no markdown fences and no extra commentary. Keep questions short, clear, and focused on the most important ideas.",
        prompt=prompt,
    )
    
    try:
        if "```json" in quiz_content:
            quiz_content = quiz_content.split("```json")[1].split("```")[0].strip()
        elif "```" in quiz_content:
            quiz_content = quiz_content.split("```")[1].split("```")[0].strip()
        
        quiz_data = json.loads(quiz_content)
    except Exception:
        quiz_data = {"questions": [], "error": "Failed to parse quiz"}
    
    quiz_id = str(uuid.uuid4())
    quiz_doc = {
        "id": quiz_id,
        "user_id": user["id"],
        "document_ids": req.document_ids,
        "quiz_data": quiz_data,
        "difficulty": req.difficulty,
        "created_at": datetime.now(timezone.utc).isoformat()
    }
    
    await db.quizzes.insert_one(quiz_doc)
    return normalize_supabase_id(quiz_doc) or dict(quiz_doc)

# Flashcards
@api_router.post("/ai/flashcards")
async def generate_flashcards(req: FlashcardRequest, request: Request):
    user = await get_current_user(request)
    doc = await db.documents.find_one({"_id": req.document_id, "user_id": user["id"], "is_deleted": False})
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    
    text = doc.get("extracted_text", "")[:5000]
    
    prompt = f"""Generate flashcards from this study material. Create 10-15 flashcards covering key concepts, definitions, and important facts.

Material:
{text}

Format as JSON:
{{
  "flashcards": [
    {{
      "front": "Question or term",
      "back": "Answer or definition",
      "hint": "Optional hint"
    }}
  ]
}}"""
    
    flashcard_content = await llm_generate(
        system="You are an expert at creating effective flashcards for studying. Always return valid JSON.",
        prompt=prompt,
    )
    
    try:
        if "```json" in flashcard_content:
            flashcard_content = flashcard_content.split("```json")[1].split("```")[0].strip()
        elif "```" in flashcard_content:
            flashcard_content = flashcard_content.split("```")[1].split("```")[0].strip()

        flashcard_data = json.loads(flashcard_content)
    except Exception:
        flashcard_data = {"flashcards": [], "error": "Failed to parse flashcards"}

    flashcard_set_id = str(uuid.uuid4())
    flashcard_doc = {
        "id": flashcard_set_id,
        "user_id": user["id"],
        "document_id": req.document_id,
        "document_name": doc.get("filename") or doc.get("name") or "Untitled Document",
        "flashcards": flashcard_data.get("flashcards", []),
        "created_at": datetime.now(timezone.utc).isoformat()
    }

    await db.flashcard_sets.insert_one(flashcard_doc)
    return normalize_supabase_id(flashcard_doc) or dict(flashcard_doc)

@api_router.get("/flashcards")
async def list_flashcards(request: Request):
    user = await get_current_user(request)
    flashcards = await db.flashcard_sets.find({"user_id": user["id"]}).to_list(1000)
    normalized_flashcards: List[Dict[str, Any]] = []
    for f in flashcards:
        normalized = ensure_response_id(f)
        if normalized is not None:
            normalized_flashcards.append(normalized)
    return normalized_flashcards

# Study planner
@api_router.post("/ai/study-plan")
async def generate_study_plan(req: StudyPlanRequest, request: Request):
    user = await get_current_user(request)
    
    user_profile = await db.users.find_one({"_id": user["id"]}, {"_id": 0}) or {}
    subjects = user_profile.get("subjects", [])
    curriculum = user_profile.get("curriculum", "")
    
    prompt = f"""Create a personalized study plan for a student.

Profile:
- Subjects: {', '.join(subjects)}
- Curriculum: {curriculum}
- Daily study time available: {req.daily_study_time} minutes
- Exam dates: {json.dumps(req.exam_dates)}

Create a weekly study schedule. Format as JSON:
{{
  "weeks": [
    {{
      "week_number": 1,
      "days": [
        {{
          "day": "Monday",
          "tasks": [
            {{
              "subject": "Math",
              "topic": "Algebra",
              "duration_minutes": 30,
              "task_type": "review"
            }}
          ]
        }}
      ]
    }}
  ]
}}"""
    
    plan_content = await llm_generate(
        system="You are an expert study planner. Create realistic, balanced study schedules. Always return valid JSON.",
        prompt=prompt,
    )
    
    try:
        if "```json" in plan_content:
            plan_content = plan_content.split("```json")[1].split("```")[0].strip()
        elif "```" in plan_content:
            plan_content = plan_content.split("```")[1].split("```")[0].strip()
        
        plan_data = json.loads(plan_content)
    except Exception:
        plan_data = {"weeks": [], "error": "Failed to parse study plan"}
    
    plan_id = str(uuid.uuid4())
    plan_doc = {
        "id": plan_id,
        "user_id": user["id"],
        "plan_data": plan_data,
        "created_at": datetime.now(timezone.utc).isoformat()
    }
    
    await db.study_plans.insert_one(plan_doc)
    return normalize_supabase_id(plan_doc) or dict(plan_doc)

@api_router.get("/study-plan/current")
async def get_current_study_plan(request: Request):
    user = await get_current_user(request)
    plan = await db.study_plans.find_one({"user_id": user["id"]}, sort=[("created_at", -1)])
    if not plan:
        return {"plan_data": {"weeks": []}}
    plan = ensure_response_id(plan) or plan
    return plan

# Progress tracking
@api_router.get("/progress/stats")
async def get_progress_stats(request: Request):
    user = await get_current_user(request)
    
    doc_count = await db.documents.count_documents({"user_id": user["id"], "is_deleted": False})
    quiz_count = await db.quizzes.count_documents({"user_id": user["id"]})
    flashcard_count = await db.flashcard_sets.count_documents({"user_id": user["id"]})
    
    streak = await db.streaks.find_one({"user_id": user["id"]}, {"_id": 0})
    if not streak:
        streak = {"current_streak": 0, "longest_streak": 0, "total_study_minutes": 0}
    
    return {
        "documents_uploaded": doc_count,
        "quizzes_taken": quiz_count,
        "flashcards_created": flashcard_count,
        "current_streak": streak.get("current_streak", 0),
        "longest_streak": streak.get("longest_streak", 0),
        "total_study_minutes": streak.get("total_study_minutes", 0)
    }

@api_router.post("/progress/log-activity")
async def log_activity(request: Request, activity_data: Dict[str, Any]):
    user = await get_current_user(request)
    
    today = datetime.now(timezone.utc).date().isoformat()
    current_streak = 1
    
    streak = await db.streaks.find_one({"user_id": user["id"]})
    if streak:
        last_date = streak.get("last_activity_date")
        current_streak = streak.get("current_streak", 0)
        longest_streak = streak.get("longest_streak", 0)
        total_minutes = streak.get("total_study_minutes", 0)
        
        if last_date != today:
            if last_date == (datetime.now(timezone.utc).date() - timedelta(days=1)).isoformat():
                current_streak += 1
            else:
                current_streak = 1
            
            longest_streak = max(longest_streak, current_streak)
            
            await db.streaks.update_one(
                {"user_id": user["id"]},
                {"$set": {
                    "current_streak": current_streak,
                    "longest_streak": longest_streak,
                    "last_activity_date": today,
                    "total_study_minutes": total_minutes + activity_data.get("study_minutes", 0)
                }}
            )
    else:
        streak_doc = {
            "id": str(uuid.uuid4()),
            "user_id": user["id"],
            "current_streak": 1,
            "longest_streak": 1,
            "last_activity_date": today,
            "total_study_minutes": activity_data.get("study_minutes", 0)
        }
        await db.streaks.insert_one(streak_doc)
    
    return {"message": "Activity logged", "current_streak": current_streak}

# Notes endpoints
@api_router.post("/notes")
async def create_note(req: NoteRequest, request: Request):
    user = await get_current_user(request)
    note_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc).isoformat()
    note_doc = {
        "id": note_id,
        "user_id": user["id"],
        "title": req.title,
        "content": req.content,
        "tags": req.tags,
        "color": req.color,
        "is_pinned": req.is_pinned,
        "document_id": req.document_id,
        "created_at": now,
        "updated_at": now,
    }
    await db.notes.insert_one(note_doc)
    return normalize_supabase_id(note_doc) or dict(note_doc)

@api_router.get("/notes")
async def list_notes(request: Request):
    user = await get_current_user(request)
    notes = await db.notes.find({"user_id": user["id"]}).sort([("is_pinned", -1), ("updated_at", -1)]).to_list(1000)
    normalized_notes: List[Dict[str, Any]] = []
    for n in notes:
        normalized = ensure_response_id(n)
        if normalized is not None:
            normalized_notes.append(normalized)
    return normalized_notes

@api_router.get("/notes/{note_id}")
async def get_note(note_id: str, request: Request):
    user = await get_current_user(request)
    note = await db.notes.find_one({"_id": note_id, "user_id": user["id"]})
    if not note:
        raise HTTPException(status_code=404, detail="Note not found")
    note = ensure_response_id(note) or note
    return note

@api_router.put("/notes/{note_id}")
async def update_note(note_id: str, req: NoteRequest, request: Request):
    user = await get_current_user(request)
    result = await db.notes.update_one(
        {"_id": note_id, "user_id": user["id"]},
        {"$set": {
            "title": req.title,
            "content": req.content,
            "tags": req.tags,
            "color": req.color,
            "is_pinned": req.is_pinned,
            "document_id": req.document_id,
            "updated_at": datetime.now(timezone.utc).isoformat(),
        }}
    )
    if result.matched_count == 0:
        raise HTTPException(status_code=404, detail="Note not found")
    note = await db.notes.find_one({"_id": note_id})
    if not note:
        raise HTTPException(status_code=404, detail="Note not found")
    note = ensure_response_id(note) or note
    return note

@api_router.delete("/notes/{note_id}")
async def delete_note(note_id: str, request: Request):
    user = await get_current_user(request)
    result = await db.notes.delete_one({"_id": note_id, "user_id": user["id"]})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Note not found")
    return {"message": "Note deleted"}

@api_router.get("/dashboard")
async def get_dashboard(request: Request):
    user = await get_current_user(request)
    
    recent_docs = await db.documents.find(
        {"user_id": user["id"], "is_deleted": False},
        {"extracted_text": 0}
    ).sort("created_at", -1).limit(5).to_list(5)
    
    normalized_recent_docs: List[Dict[str, Any]] = []
    for d in recent_docs:
        normalized = ensure_response_id(d)
        if normalized is not None:
            normalized_recent_docs.append(normalized)
    recent_docs = normalized_recent_docs

    stats = await get_progress_stats(request)
    
    study_plan = await db.study_plans.find_one({"user_id": user["id"]}, sort=[("created_at", -1)])
    
    return {
        "recent_documents": recent_docs,
        "stats": stats,
        "study_plan": study_plan.get("plan_data", {"weeks": []}) if study_plan else {"weeks": []}
    }

app.include_router(api_router)

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=os.environ.get('CORS_ORIGINS', '*').split(','),
    allow_methods=["*"],
    allow_headers=["*"],
)

async def seed_admin():
    admin_email = os.environ.get("ADMIN_EMAIL", "admin@studyai.com")
    admin_password = os.environ.get("ADMIN_PASSWORD", "admin123")
    existing = await db.users.find_one({"email": admin_email})
    if existing is None:
        admin_id = str(uuid.uuid4())
        hashed = hash_password(admin_password)
        await db.users.insert_one({
            "id": admin_id,
            "email": admin_email,
            "password_hash": hashed,
            "name": "Admin",
            "role": "admin",
            "onboarded": True,
            "created_at": datetime.now(timezone.utc).isoformat()
        })
        logger.info(f"Admin user created: {admin_email}")
    elif not verify_password(admin_password, existing["password_hash"]):
        await db.users.update_one({"email": admin_email}, {"$set": {"password_hash": hash_password(admin_password)}})
        logger.info(f"Admin password updated")
    
    # Write test credentials to the repository-local memory directory so startup
    # succeeds on systems where /app is not writable.
    mem_dir = ROOT_DIR / "memory"
    mem_dir.mkdir(parents=True, exist_ok=True)
    with open(mem_dir / "test_credentials.md", "w") as f:
        f.write(f"""# Test Credentials

## Admin Account
- Email: {admin_email}
- Password: {admin_password}
- Role: admin

## Endpoints
- POST /api/auth/register
- POST /api/auth/login
- GET /api/auth/me
- POST /api/auth/logout
- POST /api/auth/onboarding
- POST /api/documents/upload
- GET /api/documents
- POST /api/ai/summary
- POST /api/ai/tutor
- POST /api/ai/quiz
- POST /api/ai/flashcards
- POST /api/ai/study-plan
- GET /api/progress/stats
- GET /api/dashboard
""")

@app.on_event("startup")
async def startup():
    logger.info(f"Local storage directory: {STORAGE_DIR}")
    logger.info(f"Using OpenRouter model: {OPENAI_MODEL}")


    await seed_admin()

    await db.users.create_index("email", unique=True)
    await db.documents.create_index("user_id")
    await db.summaries.create_index("document_id")
    await db.quizzes.create_index("user_id")
    await db.flashcard_sets.create_index("user_id")
    await db.study_plans.create_index("user_id")
    await db.streaks.create_index("user_id")
    await db.notes.create_index("user_id")

    logger.info("Database indexes created")

@app.on_event("shutdown")
async def shutdown_db_client():
    return None
