from fastapi import FastAPI, APIRouter, File, UploadFile, Request, HTTPException, Response, Header, Query
from fastapi.responses import StreamingResponse
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import logging
import certifi
from urllib.parse import urlparse, quote_plus
from pathlib import Path
from pydantic import BaseModel, Field, EmailStr
from typing import List, Optional, Dict, Any
import uuid
from datetime import datetime, timezone, timedelta
import bcrypt
import jwt
import requests
import io
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from openai import AsyncOpenAI
import json

ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

# Load and sanitize MONGO_URL from environment. This will:
# - remove accidental angle-bracket placeholders like '<password>'
# - URL-encode the password portion so special chars don't break parsing
def _sanitize_mongo_uri(uri: str) -> str:
    if not uri:
        return uri
    try:
        p = urlparse(uri)
        if p.username or p.password:
            # remove angle brackets if present and URL-encode the password
            user = (p.username or "").replace('<', '').replace('>', '')
            pwd = (p.password or "").replace('<', '').replace('>', '')
            pwd_enc = quote_plus(pwd)
            # Safely replace only the credential portion in the original URI
            original_cred = f"{p.username}:{p.password}@"
            replacement_cred = f"{user}:{pwd_enc}@"
            if original_cred in uri:
                return uri.replace(original_cred, replacement_cred, 1)
            # fallback: return constructed minimal form
            netloc = p.netloc.split('@')[-1]
            rebuilt = f"{p.scheme}://{user}:{pwd_enc}@{netloc}"
            if p.path:
                rebuilt += p.path
            if p.query:
                rebuilt += f"?{p.query}"
            if p.fragment:
                rebuilt += f"#{p.fragment}"
            return rebuilt
    except Exception:
        # on any parsing issue, return original uri (no sensitive logging)
        return uri
    return uri

mongo_url = _sanitize_mongo_uri(os.environ.get('MONGO_URL', ''))
# Ensure driver uses system CA bundle from certifi for Atlas TLS verification
client = AsyncIOMotorClient(mongo_url, tls=True, tlsCAFile=certifi.where(), serverSelectionTimeoutMS=30000)
db = client[os.environ['DB_NAME']]

app = FastAPI()
api_router = APIRouter(prefix="/api")

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# Storage setup (local filesystem)
STORAGE_DIR = Path(os.environ.get("STORAGE_DIR", str(ROOT_DIR / "storage")))
STORAGE_DIR.mkdir(parents=True, exist_ok=True)
APP_NAME = "ai-study-assistant"

# Ollama configuration
OPENAI_API_KEY = "ollama"

OLLAMA_MODEL = os.environ.get("OLLAMA_MODEL", "gemma4:12b")
OPENAI_MODEL = OLLAMA_MODEL

openai_client = AsyncOpenAI(
    base_url="http://localhost:11434/v1",
    api_key="ollama",
)


def put_object(path: str, data: bytes, content_type: str) -> dict:
    """Save file to local filesystem storage."""
    full_path = STORAGE_DIR / path
    full_path.parent.mkdir(parents=True, exist_ok=True)
    full_path.write_bytes(data)
    return {"path": path, "size": len(data), "content_type": content_type}


def get_object(path: str) -> tuple:
    """Read file from local filesystem storage."""
    full_path = STORAGE_DIR / path
    if not full_path.exists():
        raise FileNotFoundError(f"Object not found: {path}")
    return full_path.read_bytes(), "application/octet-stream"


async def llm_generate(system: str, prompt: str, model: Optional[str] = None) -> str:
    """Generate a complete LLM response (non-streaming)."""
    if not openai_client:
        raise HTTPException(
    status_code=503,
    detail="Ollama is not running or not configured."
)
    resp = await openai_client.chat.completions.create(
        model=model or OLLAMA_MODEL,
        messages=[
            {"role": "system", "content": system},
            {"role": "user", "content": prompt},
        ],
    )
    return resp.choices[0].message.content or ""


async def llm_stream(system: str, prompt: str, model: Optional[str] = None):
    """Stream LLM response as async iterator of text chunks."""
    if not openai_client:
        raise HTTPException(
            status_code=503,
            detail="Ollama is not running or not configured.",
        )

    stream = await openai_client.chat.completions.create(
        model=model or OLLAMA_MODEL,
        messages=[
            {"role": "system", "content": system},
            {"role": "user", "content": prompt},
        ],
        stream=True,
    )

    async for chunk in stream:
        delta = chunk.choices[0].delta.content if chunk.choices else None
        if delta:
            yield delta
# Auth utilities
JWT_ALGORITHM = "HS256"

def get_jwt_secret() -> str:
    return os.environ["JWT_SECRET"]

def hash_password(password: str) -> str:
    salt = bcrypt.gensalt()
    hashed = bcrypt.hashpw(password.encode("utf-8"), salt)
    return hashed.decode("utf-8")

def verify_password(plain_password: str, hashed_password: str) -> bool:
    return bcrypt.checkpw(plain_password.encode("utf-8"), hashed_password.encode("utf-8"))

def create_access_token(user_id: str, email: str) -> str:
    payload = {"sub": user_id, "email": email, "exp": datetime.now(timezone.utc) + timedelta(minutes=15), "type": "access"}
    return jwt.encode(payload, get_jwt_secret(), algorithm=JWT_ALGORITHM)

def create_refresh_token(user_id: str) -> str:
    payload = {"sub": user_id, "exp": datetime.now(timezone.utc) + timedelta(days=7), "type": "refresh"}
    return jwt.encode(payload, get_jwt_secret(), algorithm=JWT_ALGORITHM)

async def get_current_user(request: Request) -> dict:
    token = request.cookies.get("access_token")
    if not token:
        auth_header = request.headers.get("Authorization", "")
        if auth_header.startswith("Bearer "):
            token = auth_header[7:]
    if not token:
        raise HTTPException(status_code=401, detail="Not authenticated")
    try:
        payload = jwt.decode(token, get_jwt_secret(), algorithms=[JWT_ALGORITHM])
        if payload.get("type") != "access":
            raise HTTPException(status_code=401, detail="Invalid token type")
        user = await db.users.find_one({"_id": payload["sub"]})
        if not user:
            raise HTTPException(status_code=401, detail="User not found")
        user["id"] = user.pop("_id")
        user.pop("password_hash", None)
        return user
    except jwt.ExpiredSignatureError:
        raise HTTPException(status_code=401, detail="Token expired")
    except jwt.InvalidTokenError:
        raise HTTPException(status_code=401, detail="Invalid token")

# Models
class RegisterRequest(BaseModel):
    email: EmailStr
    password: str
    name: str

class LoginRequest(BaseModel):
    email: EmailStr
    password: str

class OnboardingRequest(BaseModel):
    grade_level: str
    subjects: List[str]
    curriculum: str
    exam_dates: Optional[List[Dict[str, str]]] = []
    daily_study_goal: int = 60

class SummaryRequest(BaseModel):
    document_id: str
    summary_type: str = "detailed"

class TutorRequest(BaseModel):
    document_id: str
    question: str
    difficulty_mode: str = "intermediate"
    session_id: Optional[str] = None

class QuizRequest(BaseModel):
    document_ids: List[str]
    question_types: List[str]
    difficulty: str = "medium"
    num_questions: int = 10

class FlashcardRequest(BaseModel):
    document_id: str

class StudyPlanRequest(BaseModel):
    exam_dates: List[Dict[str, str]]
    daily_study_time: int

class NoteRequest(BaseModel):
    title: str
    content: str = ""
    tags: List[str] = []
    color: str = "butter"
    is_pinned: bool = False
    document_id: Optional[str] = None

# Helper: Extract text from files
def extract_text_from_file(file_content: bytes, file_type: str) -> str:
    try:
        if file_type == "application/pdf":
            pdf_reader = PdfReader(io.BytesIO(file_content))
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text
        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = DocxDocument(io.BytesIO(file_content))
            text = "\n".join([para.text for para in doc.paragraphs])
            return text
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            prs = Presentation(io.BytesIO(file_content))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    txt = getattr(shape, "text", None)
                    if txt:
                        text += txt + "\n"
            return text
        elif file_type == "text/plain":
            return file_content.decode("utf-8")
        else:
            return ""
    except Exception as e:
        logger.error(f"Error extracting text: {e}")
        return ""

# Auth endpoints
@api_router.post("/auth/register")
async def register(req: RegisterRequest, response: Response):
    email = req.email.lower()
    existing = await db.users.find_one({"email": email})
    if existing:
        raise HTTPException(status_code=400, detail="Email already registered")
    
    user_id = str(uuid.uuid4())
    password_hash = hash_password(req.password)
    
    user_doc = {
        "_id": user_id,
        "email": email,
        "name": req.name,
        "password_hash": password_hash,
        "onboarded": False,
        "created_at": datetime.now(timezone.utc).isoformat()
    }
    
    await db.users.insert_one(user_doc)
    
    access_token = create_access_token(user_id, email)
    refresh_token = create_refresh_token(user_id)
    
    response.set_cookie(key="access_token", value=access_token, httponly=True, secure=False, samesite="lax", max_age=900, path="/")
    response.set_cookie(key="refresh_token", value=refresh_token, httponly=True, secure=False, samesite="lax", max_age=604800, path="/")
    
    user_doc.pop("password_hash")
    user_doc["id"] = user_doc.pop("_id")
    return user_doc

@api_router.post("/auth/login")
async def login(req: LoginRequest, response: Response):
    email = req.email.lower()
    user = await db.users.find_one({"email": email})
    if not user or not verify_password(req.password, user["password_hash"]):
        raise HTTPException(status_code=401, detail="Invalid email or password")
    
    user_id = str(user["_id"])
    access_token = create_access_token(user_id, email)
    refresh_token = create_refresh_token(user_id)
    
    response.set_cookie(key="access_token", value=access_token, httponly=True, secure=False, samesite="lax", max_age=900, path="/")
    response.set_cookie(key="refresh_token", value=refresh_token, httponly=True, secure=False, samesite="lax", max_age=604800, path="/")
    
    user.pop("password_hash")
    user["id"] = str(user.pop("_id"))
    return user

@api_router.post("/auth/logout")
async def logout(response: Response):
    response.delete_cookie(key="access_token", path="/")
    response.delete_cookie(key="refresh_token", path="/")
    return {"message": "Logged out successfully"}

@api_router.get("/auth/me")
async def get_me(request: Request):
    user = await get_current_user(request)
    return user

@api_router.post("/auth/onboarding")
async def complete_onboarding(req: OnboardingRequest, request: Request):
    user = await get_current_user(request)
    
    await db.users.update_one(
        {"_id": user["id"]},
        {"$set": {
            "grade_level": req.grade_level,
            "subjects": req.subjects,
            "curriculum": req.curriculum,
            "exam_dates": req.exam_dates,
            "daily_study_goal": req.daily_study_goal,
            "onboarded": True
        }}
    )
    
    streak_doc = {
        "_id": str(uuid.uuid4()),
        "user_id": user["id"],
        "current_streak": 0,
        "longest_streak": 0,
        "last_activity_date": None,
        "total_study_minutes": 0
    }
    await db.streaks.insert_one(streak_doc)
    
    return {"message": "Onboarding completed"}

# File upload
@api_router.post("/documents/upload")
async def upload_document(request: Request, file: UploadFile = File(...)):
    user = await get_current_user(request)
    
    file_content = await file.read()
    filename = getattr(file, "filename", "") or ""
    ext = filename.split(".")[-1] if "." in filename else "bin"
    storage_path = f"{APP_NAME}/uploads/{user['id']}/{uuid.uuid4()}.{ext}"
    
    content_type = file.content_type or "application/octet-stream"
    result = put_object(storage_path, file_content, content_type)
    
    extracted_text = extract_text_from_file(file_content, content_type)
    
    doc_id = str(uuid.uuid4())
    doc = {
        "_id": doc_id,
        "user_id": user["id"],
        "original_filename": file.filename,
        "storage_path": result["path"],
        "content_type": content_type,
        "size": result["size"],
        "extracted_text": extracted_text,
        "is_deleted": False,
        "created_at": datetime.now(timezone.utc).isoformat()
    }
    
    await db.documents.insert_one(doc)
    
    doc["id"] = doc.pop("_id")
    doc.pop("extracted_text", None)
    return doc

@api_router.get("/documents")
async def list_documents(request: Request):
    user = await get_current_user(request)
    docs = await db.documents.find({"user_id": user["id"], "is_deleted": False}, {"extracted_text": 0}).to_list(1000)
    for d in docs:
        d["id"] = d.pop("_id")
    return docs

@api_router.get("/documents/{doc_id}")
async def get_document(doc_id: str, request: Request):
    user = await get_current_user(request)
    doc = await db.documents.find_one({"_id": doc_id, "user_id": user["id"], "is_deleted": False})
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    doc["id"] = doc.pop("_id")
    return doc

@api_router.delete("/documents/{doc_id}")
async def delete_document(doc_id: str, request: Request):
    user = await get_current_user(request)
    result = await db.documents.update_one({"_id": doc_id, "user_id": user["id"]}, {"$set": {"is_deleted": True}})
    if result.modified_count == 0:
        raise HTTPException(status_code=404, detail="Document not found")
    return {"message": "Document deleted"}

# AI Summary
@api_router.post("/ai/summary")
async def generate_summary(req: SummaryRequest, request: Request):
    user = await get_current_user(request)
    doc = await db.documents.find_one({"_id": req.document_id, "user_id": user["id"], "is_deleted": False})
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    
    existing = await db.summaries.find_one({"document_id": req.document_id, "summary_type": req.summary_type})
    if existing:
        existing["id"] = str(existing.pop("_id"))
        return existing
    
    text = doc.get("extracted_text", "")
    if not text:
        raise HTTPException(status_code=400, detail="No text content found in document")
    
    prompt_templates = {
        "short": f"Provide a concise summary (3-5 sentences) of the following study material:\n\n{text[:4000]}",
        "detailed": f"Create a comprehensive summary with:\n1. Main Topics\n2. Key Concepts\n3. Important Definitions\n4. Formulas (if any)\n5. Exam Tips\n\nMaterial:\n{text[:6000]}",
        "bullet_points": f"Summarize in clear bullet points:\n\n{text[:5000]}",
        "timeline": f"If applicable, create a timeline of events. Otherwise, organize by sequence:\n\n{text[:5000]}"
    }
    
    prompt = prompt_templates.get(req.summary_type, prompt_templates["detailed"])
    
    summary_content = await llm_generate(
        system="You are an expert study assistant. Create clear, well-organized summaries for students.",
        prompt=prompt,
    )
    
    summary_id = str(uuid.uuid4())
    summary_doc = {
        "_id": summary_id,
        "document_id": req.document_id,
        "user_id": user["id"],
        "summary_type": req.summary_type,
        "content": summary_content,
        "created_at": datetime.now(timezone.utc).isoformat()
    }
    
    await db.summaries.insert_one(summary_doc)
    summary_doc["id"] = summary_doc.pop("_id")
    return summary_doc

# AI Tutor (streaming)
@api_router.post("/ai/tutor")
async def ai_tutor(req: TutorRequest, request: Request):
    user = await get_current_user(request)
    doc = await db.documents.find_one({"_id": req.document_id, "user_id": user["id"], "is_deleted": False})
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    
    text = doc.get("extracted_text", "")[:6000]
    
    mode_prompts = {
        "beginner": "Explain concepts in very simple terms, use analogies, and avoid jargon.",
        "intermediate": "Provide clear explanations with examples. Balance simplicity with accuracy.",
        "advanced": "Provide in-depth explanations. You can use technical terminology.",
        "exam_prep": "Focus on exam-relevant information, common pitfalls, and test strategies."
    }
    
    system_msg = f"{mode_prompts.get(req.difficulty_mode, mode_prompts['intermediate'])} You are a patient tutor. Guide students with hints before giving direct answers. Reference this material: {text}"
    
    async def event_generator():
        async for delta in llm_stream(system=system_msg, prompt=req.question):
            yield f"data: {json.dumps({'content': delta})}\n\n"
        yield f"data: {json.dumps({'done': True})}\n\n"
    
    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"}
    )

# Quiz generation
@api_router.post("/ai/quiz")
async def generate_quiz(req: QuizRequest, request: Request):
    user = await get_current_user(request)
    
    combined_text = ""
    for doc_id in req.document_ids:
        doc = await db.documents.find_one({"_id": doc_id, "user_id": user["id"], "is_deleted": False})
        if doc:
            combined_text += doc.get("extracted_text", "")[:2000] + "\n\n"
    
    if not combined_text:
        raise HTTPException(status_code=400, detail="No text found in documents")
    
    question_types_str = ", ".join(req.question_types)
    
    prompt = f"""Generate a quiz with {req.num_questions} questions from this material.
Question types: {question_types_str}
Difficulty: {req.difficulty}

Material:
{combined_text[:5000]}

Format as JSON:
{{
  "questions": [
    {{
      "type": "multiple_choice",
      "question": "...",
      "options": ["A", "B", "C", "D"],
      "correct_answer": "A",
      "explanation": "..."
    }}
  ]
}}"""
    
    quiz_content = await llm_generate(
        system="You are an expert quiz generator. Always return valid JSON.",
        prompt=prompt,
    )
    
    try:
        if "```json" in quiz_content:
            quiz_content = quiz_content.split("```json")[1].split("```")[0].strip()
        elif "```" in quiz_content:
            quiz_content = quiz_content.split("```")[1].split("```")[0].strip()
        
        quiz_data = json.loads(quiz_content)
    except Exception:
        quiz_data = {"questions": [], "error": "Failed to parse quiz"}
    
    quiz_id = str(uuid.uuid4())
    quiz_doc = {
        "_id": quiz_id,
        "user_id": user["id"],
        "document_ids": req.document_ids,
        "quiz_data": quiz_data,
        "difficulty": req.difficulty,
        "created_at": datetime.now(timezone.utc).isoformat()
    }
    
    await db.quizzes.insert_one(quiz_doc)
    quiz_doc["id"] = quiz_doc.pop("_id")
    return quiz_doc

# Flashcards
@api_router.post("/ai/flashcards")
async def generate_flashcards(req: FlashcardRequest, request: Request):
    user = await get_current_user(request)
    doc = await db.documents.find_one({"_id": req.document_id, "user_id": user["id"], "is_deleted": False})
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    
    text = doc.get("extracted_text", "")[:5000]
    
    prompt = f"""Generate flashcards from this study material. Create 10-15 flashcards covering key concepts, definitions, and important facts.

Material:
{text}

Format as JSON:
{{
  "flashcards": [
    {{
      "front": "Question or term",
      "back": "Answer or definition",
      "hint": "Optional hint"
    }}
  ]
}}"""
    
    flashcard_content = await llm_generate(
        system="You are an expert at creating effective flashcards for studying. Always return valid JSON.",
        prompt=prompt,
    )
    
    try:
        if "```json" in flashcard_content:
            flashcard_content = flashcard_content.split("```json")[1].split("```")[0].strip()
        elif "```" in flashcard_content:
            flashcard_content = flashcard_content.split("```")[1].split("```")[0].strip()
        
        flashcard_data = json.loads(flashcard_content)
    except Exception:
        flashcard_data = {"flashcards": [], "error": "Failed to parse flashcards"}
    
    flashcard_set_id = str(uuid.uuid4())
    flashcard_doc = {
        "_id": flashcard_set_id,
        "user_id": user["id"],
        "document_id": req.document_id,
        "flashcards": flashcard_data.get("flashcards", []),
        "created_at": datetime.now(timezone.utc).isoformat()
    }
    
    await db.flashcard_sets.insert_one(flashcard_doc)
    flashcard_doc["id"] = flashcard_doc.pop("_id")
    return flashcard_doc

@api_router.get("/flashcards")
async def list_flashcards(request: Request):
    user = await get_current_user(request)
    flashcards = await db.flashcard_sets.find({"user_id": user["id"]}).to_list(1000)
    for f in flashcards:
        f["id"] = f.pop("_id")
    return flashcards

# Study planner
@api_router.post("/ai/study-plan")
async def generate_study_plan(req: StudyPlanRequest, request: Request):
    user = await get_current_user(request)
    
    user_profile = await db.users.find_one({"_id": user["id"]}, {"_id": 0}) or {}
    subjects = user_profile.get("subjects", [])
    curriculum = user_profile.get("curriculum", "")
    
    prompt = f"""Create a personalized study plan for a student.

Profile:
- Subjects: {', '.join(subjects)}
- Curriculum: {curriculum}
- Daily study time available: {req.daily_study_time} minutes
- Exam dates: {json.dumps(req.exam_dates)}

Create a weekly study schedule. Format as JSON:
{{
  "weeks": [
    {{
      "week_number": 1,
      "days": [
        {{
          "day": "Monday",
          "tasks": [
            {{
              "subject": "Math",
              "topic": "Algebra",
              "duration_minutes": 30,
              "task_type": "review"
            }}
          ]
        }}
      ]
    }}
  ]
}}"""
    
    plan_content = await llm_generate(
        system="You are an expert study planner. Create realistic, balanced study schedules. Always return valid JSON.",
        prompt=prompt,
    )
    
    try:
        if "```json" in plan_content:
            plan_content = plan_content.split("```json")[1].split("```")[0].strip()
        elif "```" in plan_content:
            plan_content = plan_content.split("```")[1].split("```")[0].strip()
        
        plan_data = json.loads(plan_content)
    except Exception:
        plan_data = {"weeks": [], "error": "Failed to parse study plan"}
    
    plan_id = str(uuid.uuid4())
    plan_doc = {
        "_id": plan_id,
        "user_id": user["id"],
        "plan_data": plan_data,
        "created_at": datetime.now(timezone.utc).isoformat()
    }
    
    await db.study_plans.insert_one(plan_doc)
    plan_doc["id"] = plan_doc.pop("_id")
    return plan_doc

@api_router.get("/study-plan/current")
async def get_current_study_plan(request: Request):
    user = await get_current_user(request)
    plan = await db.study_plans.find_one({"user_id": user["id"]}, sort=[("created_at", -1)])
    if not plan:
        return {"plan_data": {"weeks": []}}
    plan["id"] = plan.pop("_id")
    return plan

# Progress tracking
@api_router.get("/progress/stats")
async def get_progress_stats(request: Request):
    user = await get_current_user(request)
    
    doc_count = await db.documents.count_documents({"user_id": user["id"], "is_deleted": False})
    quiz_count = await db.quizzes.count_documents({"user_id": user["id"]})
    flashcard_count = await db.flashcard_sets.count_documents({"user_id": user["id"]})
    
    streak = await db.streaks.find_one({"user_id": user["id"]}, {"_id": 0})
    if not streak:
        streak = {"current_streak": 0, "longest_streak": 0, "total_study_minutes": 0}
    
    return {
        "documents_uploaded": doc_count,
        "quizzes_taken": quiz_count,
        "flashcards_created": flashcard_count,
        "current_streak": streak.get("current_streak", 0),
        "longest_streak": streak.get("longest_streak", 0),
        "total_study_minutes": streak.get("total_study_minutes", 0)
    }

@api_router.post("/progress/log-activity")
async def log_activity(request: Request, activity_data: Dict[str, Any]):
    user = await get_current_user(request)
    
    today = datetime.now(timezone.utc).date().isoformat()
    current_streak = 1
    
    streak = await db.streaks.find_one({"user_id": user["id"]})
    if streak:
        last_date = streak.get("last_activity_date")
        current_streak = streak.get("current_streak", 0)
        longest_streak = streak.get("longest_streak", 0)
        total_minutes = streak.get("total_study_minutes", 0)
        
        if last_date != today:
            if last_date == (datetime.now(timezone.utc).date() - timedelta(days=1)).isoformat():
                current_streak += 1
            else:
                current_streak = 1
            
            longest_streak = max(longest_streak, current_streak)
            
            await db.streaks.update_one(
                {"user_id": user["id"]},
                {"$set": {
                    "current_streak": current_streak,
                    "longest_streak": longest_streak,
                    "last_activity_date": today,
                    "total_study_minutes": total_minutes + activity_data.get("study_minutes", 0)
                }}
            )
    else:
        streak_doc = {
            "_id": str(uuid.uuid4()),
            "user_id": user["id"],
            "current_streak": 1,
            "longest_streak": 1,
            "last_activity_date": today,
            "total_study_minutes": activity_data.get("study_minutes", 0)
        }
        await db.streaks.insert_one(streak_doc)
    
    return {"message": "Activity logged", "current_streak": current_streak}

# Notes endpoints
@api_router.post("/notes")
async def create_note(req: NoteRequest, request: Request):
    user = await get_current_user(request)
    note_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc).isoformat()
    note_doc = {
        "_id": note_id,
        "user_id": user["id"],
        "title": req.title,
        "content": req.content,
        "tags": req.tags,
        "color": req.color,
        "is_pinned": req.is_pinned,
        "document_id": req.document_id,
        "created_at": now,
        "updated_at": now,
    }
    await db.notes.insert_one(note_doc)
    note_doc["id"] = note_doc.pop("_id")
    return note_doc

@api_router.get("/notes")
async def list_notes(request: Request):
    user = await get_current_user(request)
    notes = await db.notes.find({"user_id": user["id"]}).sort([("is_pinned", -1), ("updated_at", -1)]).to_list(1000)
    for n in notes:
        n["id"] = n.pop("_id")
    return notes

@api_router.get("/notes/{note_id}")
async def get_note(note_id: str, request: Request):
    user = await get_current_user(request)
    note = await db.notes.find_one({"_id": note_id, "user_id": user["id"]})
    if not note:
        raise HTTPException(status_code=404, detail="Note not found")
    note["id"] = note.pop("_id")
    return note

@api_router.put("/notes/{note_id}")
async def update_note(note_id: str, req: NoteRequest, request: Request):
    user = await get_current_user(request)
    result = await db.notes.update_one(
        {"_id": note_id, "user_id": user["id"]},
        {"$set": {
            "title": req.title,
            "content": req.content,
            "tags": req.tags,
            "color": req.color,
            "is_pinned": req.is_pinned,
            "document_id": req.document_id,
            "updated_at": datetime.now(timezone.utc).isoformat(),
        }}
    )
    if result.matched_count == 0:
        raise HTTPException(status_code=404, detail="Note not found")
    note = await db.notes.find_one({"_id": note_id})
    if not note:
        raise HTTPException(status_code=404, detail="Note not found")
    note["id"] = note.pop("_id")
    return note

@api_router.delete("/notes/{note_id}")
async def delete_note(note_id: str, request: Request):
    user = await get_current_user(request)
    result = await db.notes.delete_one({"_id": note_id, "user_id": user["id"]})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Note not found")
    return {"message": "Note deleted"}

@api_router.get("/dashboard")
async def get_dashboard(request: Request):
    user = await get_current_user(request)
    
    recent_docs = await db.documents.find(
        {"user_id": user["id"], "is_deleted": False},
        {"extracted_text": 0}
    ).sort("created_at", -1).limit(5).to_list(5)
    
    for d in recent_docs:
        d["id"] = d.pop("_id")
    
    stats = await get_progress_stats(request)
    
    study_plan = await db.study_plans.find_one({"user_id": user["id"]}, sort=[("created_at", -1)])
    
    return {
        "recent_documents": recent_docs,
        "stats": stats,
        "study_plan": study_plan.get("plan_data", {"weeks": []}) if study_plan else {"weeks": []}
    }

app.include_router(api_router)

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=os.environ.get('CORS_ORIGINS', '*').split(','),
    allow_methods=["*"],
    allow_headers=["*"],
)

async def seed_admin():
    admin_email = os.environ.get("ADMIN_EMAIL", "admin@studyai.com")
    admin_password = os.environ.get("ADMIN_PASSWORD", "admin123")
    existing = await db.users.find_one({"email": admin_email})
    if existing is None:
        admin_id = str(uuid.uuid4())
        hashed = hash_password(admin_password)
        await db.users.insert_one({
            "_id": admin_id,
            "email": admin_email,
            "password_hash": hashed,
            "name": "Admin",
            "role": "admin",
            "onboarded": True,
            "created_at": datetime.now(timezone.utc).isoformat()
        })
        logger.info(f"Admin user created: {admin_email}")
    elif not verify_password(admin_password, existing["password_hash"]):
        await db.users.update_one({"email": admin_email}, {"$set": {"password_hash": hash_password(admin_password)}})
        logger.info(f"Admin password updated")
    
    # Write test credentials to the repository-local memory directory so startup
    # succeeds on systems where /app is not writable.
    mem_dir = ROOT_DIR / "memory"
    mem_dir.mkdir(parents=True, exist_ok=True)
    with open(mem_dir / "test_credentials.md", "w") as f:
        f.write(f"""# Test Credentials

## Admin Account
- Email: {admin_email}
- Password: {admin_password}
- Role: admin

## Endpoints
- POST /api/auth/register
- POST /api/auth/login
- GET /api/auth/me
- POST /api/auth/logout
- POST /api/auth/onboarding
- POST /api/documents/upload
- GET /api/documents
- POST /api/ai/summary
- POST /api/ai/tutor
- POST /api/ai/quiz
- POST /api/ai/flashcards
- POST /api/ai/study-plan
- GET /api/progress/stats
- GET /api/dashboard
""")

@app.on_event("startup")
async def startup():
    logger.info(f"Local storage directory: {STORAGE_DIR}")
    logger.info(f"Using Ollama model: {OLLAMA_MODEL}")

    await seed_admin()

    await db.users.create_index("email", unique=True)
    await db.documents.create_index("user_id")
    await db.summaries.create_index("document_id")
    await db.quizzes.create_index("user_id")
    await db.flashcard_sets.create_index("user_id")
    await db.study_plans.create_index("user_id")
    await db.streaks.create_index("user_id")
    await db.notes.create_index("user_id")

    logger.info("Database indexes created")

@app.on_event("shutdown")
async def shutdown_db_client():
    client.close()
